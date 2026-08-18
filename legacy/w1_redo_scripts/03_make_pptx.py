#!/usr/bin/env python3
"""Build an editable weekly-report PPTX (with speaker notes) for the tumor-only
HCC1395 ClairS-TO candidate SNV analysis. Plain, explanatory English."""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.join(os.path.dirname(__file__), "..")
OUT = os.path.join(BASE, "weekly_report_HCC1395_clairsto.pptx")

NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x2B, 0x8C, 0xBE)
GREY = RGBColor(0x55, 0x55, 0x55)
LIGHT = RGBColor(0xF2, 0xF6, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def add_slide():
    return prs.slides.add_slide(BLANK)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def band(slide):
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(1.0))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    bar.shadow.inherit = False
    return bar


def title(slide, text, sub=None):
    band(slide)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.12), SW - Inches(1), Inches(0.9))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = text
    r.font.size = Pt(26); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    if sub:
        p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = sub
        r2.font.size = Pt(12); r2.font.color.rgb = RGBColor(0xD6, 0xE4, 0xF0)


def bullets(slide, items, left=0.6, top=1.3, width=7.0, height=5.6, size=16):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    first = True
    for lvl, txt, *bold in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl
        r = p.add_run(); r.text = ("• " if lvl == 0 else "– ") + txt
        r.font.size = Pt(size - lvl * 2)
        r.font.color.rgb = NAVY if lvl == 0 else GREY
        if bold and bold[0]:
            r.font.bold = True; r.font.color.rgb = BLUE
        p.space_after = Pt(6)
    return tb


def pic(slide, path, left, top, width):
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), width=Inches(width))


def metric_table(slide, rows, left=0.6, top=1.4, width=6.2):
    r, c = len(rows), 2
    gt = slide.shapes.add_table(r, c, Inches(left), Inches(top), Inches(width), Inches(0.5 * r)).table
    gt.columns[0].width = Inches(width * 0.68); gt.columns[1].width = Inches(width * 0.32)
    for i, (k, v) in enumerate(rows):
        for j, val in enumerate((k, v)):
            cell = gt.cell(i, j); cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.runs[0].font.size = Pt(13)
            if i == 0:
                para.runs[0].font.bold = True
                para.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT if i % 2 else RGBColor(0xFF, 0xFF, 0xFF)
                if j == 1:
                    para.runs[0].font.bold = True; para.runs[0].font.color.rgb = BLUE
            if j == 1:
                para.alignment = PP_ALIGN.CENTER


# ---------- Slide 1: Title ----------
s = add_slide()
bg = s.shapes.add_shape(1, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background(); bg.shadow.inherit = False
tb = s.shapes.add_textbox(Inches(0.8), Inches(2.3), SW - Inches(1.6), Inches(2.5))
tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = "Weekly Research Report"
r.font.size = Pt(20); r.font.color.rgb = RGBColor(0xAF, 0xCB, 0xE3)
p = tf.add_paragraph(); r = p.add_run()
r.text = "Tumor-only somatic SNV candidate set from HCC1395 (ClairS-TO)"
r.font.size = Pt(32); r.font.bold = True; r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p = tf.add_paragraph(); r = p.add_run()
r.text = "First step toward tumor-informed markers for MRD dilution analysis"
r.font.size = Pt(16); r.font.color.rgb = RGBColor(0xD6, 0xE4, 0xF0)
p = tf.add_paragraph(); p.space_before = Pt(20); r = p.add_run()
r.text = "Long-read ONT  •  ClairS-TO v0.5.0  •  strictly tumor-only (no matched normal)"
r.font.size = Pt(13); r.font.color.rgb = RGBColor(0xAF, 0xCB, 0xE3)
notes(s, "This week I generated and characterized a tumor-only somatic SNV candidate set "
      "from the HCC1395 tumor long-read BAM using ClairS-TO. The purpose is to build a list "
      "of predefined mutation loci that we can later look at in the diluted MRD samples. "
      "Two things to stress up front: (1) everything is tumor-only, no matched normal was used; "
      "(2) I call these 'candidate' variants, not confirmed somatic mutations.")

# ---------- Slide 2: Objective / question ----------
s = add_slide(); title(s, "Objective & research question")
bullets(s, [
    (0, "Project context: tumor-only MRD using long-read sequencing.", True),
    (0, "This week is the FIRST step — not full MRDetect yet.", False),
    (0, "Goal: build a tumor-only somatic SNV candidate set from HCC1395 tumor reads,", True),
    (1, "so these loci can later be used as predefined tumor-informed markers", False),
    (1, "in the 1% / 0.1% / 0.01% dilution BAMs.", False),
    (0, "Research question:", True),
    (1, "Can ClairS-TO generate a usable tumor-only somatic SNV candidate set", False),
    (1, "from HCC1395 for downstream MRD dilution analysis?", False),
], width=11.8)
notes(s, "The big project is measuring minimal residual disease from long reads without a matched "
      "normal. Before doing any detection, we need a marker list. So the only question this week is "
      "narrow and practical: can the tumor-only caller ClairS-TO give us a usable set of candidate "
      "SNVs from HCC1395? Usable means: enough of them, good enough quality, and technically ready "
      "to interrogate in the dilution samples.")

# ---------- Slide 3: Approach ----------
s = add_slide(); title(s, "Approach — how the candidate set is built")
bullets(s, [
    (0, "HCC1395 tumor long-read BAM (ONT, ~100x, tumor-only)", True),
    (0, "→ ClairS-TO tumor-only somatic calling (v0.5.0, ssrs model, SNV-only)", True),
    (0, "→ keep PASS SNVs (passed quality + panel-of-normals + artifact filters)", True),
    (0, "→ basic QC: VAF, depth, per-chromosome (characterize, do not over-filter)", True),
    (0, "→ predefined loci for later dilution analysis", True),
    (1, "Reused the existing ClairS-TO run; did not modify BAM/VCF.", False),
    (1, "All steps kept as reproducible scripts in scripts/.", False),
], width=11.8)
notes(s, "The pipeline is a simple chain. Start from the tumor BAM. Run ClairS-TO in tumor-only mode. "
      "Because there is no matched normal, ClairS-TO removes germline using public population "
      "databases (a 'panel of normals') and tags obvious artifacts. We then keep only PASS SNVs. "
      "Finally we do light QC to describe the set — we deliberately do not aggressively filter, "
      "because the goal this week is characterization, not tuning. I reused an existing run rather "
      "than recomputing, and everything is scripted so it can be reproduced.")

# ---------- Slide 4: Inputs ----------
s = add_slide(); title(s, "Inputs & configuration")
metric_table(s, [
    ("Item", "Value"),
    ("Tumor BAM", "HCC1395 ONT 5kHz (tumor-only)"),
    ("Reference", "GRCh38 no_alt_analysis_set"),
    ("Platform", "ONT R10, minimap2 map-ont"),
    ("Coverage", "~85-100x"),
    ("Caller", "ClairS-TO v0.5.0 (ssrs)"),
    ("Mode / filters", "SNV-only, min AF 0.05, QUAL>4"),
    ("Germline removal", "gnomAD/dbSNP/1000g/CoLoRSdb PoN"),
], left=0.6, top=1.35, width=7.2)
bullets(s, [
    (0, "Strictly tumor-only:", True),
    (1, "HCC1395BL / matched normal NOT used anywhere.", False),
    (0, "Read-only:", True),
    (1, "original BAM and VCF were not modified.", False),
    (0, "Assumption logged:", True),
    (1, "indel calling was disabled in this run,", False),
    (1, "so PASS indels = 0 by configuration.", False),
], left=8.1, top=1.35, width=4.9, size=14)
notes(s, "These are the exact inputs so the work is traceable. The tumor sample is ONT long-read at "
      "roughly 85 to 100x, aligned to GRCh38. The caller is ClairS-TO version 0.5.0 with the ssrs "
      "model. Two configuration points matter: the minimum allele fraction is 0.05, which sets a "
      "floor on how low a VAF we can call; and indels were turned off, so there are zero PASS indels "
      "by design, not because none exist. Germline is handled purely by population databases since "
      "we have no matched normal.")

# ---------- Slide 5: Key results - numbers ----------
s = add_slide(); title(s, "Key results — candidate counts")
metric_table(s, [
    ("Metric", "Value"),
    ("Total ClairS-TO calls", "3,169,996"),
    ("PASS calls", "48,819"),
    ("PASS SNVs", "48,819"),
    ("PASS indels (disabled)", "0"),
    ("Preliminary candidate SNVs", "48,819"),
], left=0.6, top=1.5, width=6.4)
bullets(s, [
    (0, "~3.17M raw calls, but almost all are tagged", True),
    (1, "NonSomatic (germline/population) and dropped.", False),
    (0, "48,819 PASS SNVs remain as the candidate set.", True),
    (0, "This is a large, workable marker list", True),
    (1, "for a hypermutated line like HCC1395.", False),
], left=7.3, top=1.6, width=5.6, size=15)
notes(s, "The headline number is 48,819 candidate PASS SNVs. Notice the funnel: ClairS-TO looks at "
      "over three million positions, but the vast majority are tagged as non-somatic — mostly "
      "germline caught by the population databases — and removed. What survives as PASS is about "
      "49 thousand SNVs. For a heavily mutated cell line like HCC1395 that is a reasonable, usable "
      "number of markers to carry forward.")

# ---------- Slide 6: VAF ----------
s = add_slide(); title(s, "Result — VAF distribution")
pic(s, os.path.join(BASE, "vaf_distribution.png"), 0.5, 1.3, 7.2)
bullets(s, [
    (0, "Median VAF 0.46 (middle 50%: 0.29-0.78).", True),
    (0, "~24% sit near 0.5 (heterozygous-like).", True),
    (0, "~22% sit above 0.9 (homozygous-like).", True),
    (0, "~0% below VAF 0.05 — floored by min AF 0.05.", True),
    (1, "So the set is biased to clonal / higher-VAF loci.", False),
    (0, "VAF = FORMAT/AF (ClairS-TO tumor allele fraction).", False),
], left=8.0, top=1.4, width=5.0, size=14)
notes(s, "VAF is the fraction of reads supporting the alternate allele. The distribution has a broad "
      "peak with a median near 0.46. Two visible clusters: one around 0.5, which looks heterozygous, "
      "and one near 1.0, which looks homozygous. Important caution: in a tumor-only setting these "
      "germline-LIKE shapes are not proof of germline — a hypermutated tumor can have many clonal "
      "somatic SNVs at these same VAFs, and loss-of-heterozygosity can push a real somatic variant "
      "toward VAF 1.0. Also, because we set the minimum AF to 0.05, there are essentially no very "
      "low-VAF calls, so this candidate set leans toward clonal markers.")

# ---------- Slide 7: Depth ----------
s = add_slide(); title(s, "Result — depth distribution")
pic(s, os.path.join(BASE, "depth_distribution.png"), 0.5, 1.3, 7.2)
bullets(s, [
    (0, "Median depth 80x (middle 50%: 59-106x).", True),
    (0, "Only 0.2% of calls below 10x — good support.", True),
    (0, "Small 0.4% tail above ~240x (3x median).", True),
    (1, "likely collapsed repeats or CNV amplification.", False),
    (0, "depth = FORMAT/DP; alt/ref = FORMAT/AD.", False),
], left=8.0, top=1.4, width=5.0, size=14)
notes(s, "Depth is how many reads cover each candidate. The median is about 80x, which is comfortable "
      "for long-read calling, and very few calls sit below 10x, so read support is generally solid. "
      "There is a thin high-depth tail above roughly 240x — those are worth a glance later because "
      "extreme depth often means collapsed repeats or amplified copy-number regions, which can "
      "produce artifacts. Nothing here is alarming, but I flag it rather than silently drop it.")

# ---------- Slide 8: per-chrom + scatter ----------
s = add_slide(); title(s, "Result — genome distribution & VAF vs depth")
pic(s, os.path.join(BASE, "variants_per_chromosome.png"), 0.4, 1.35, 6.6)
pic(s, os.path.join(BASE, "vaf_vs_depth.png"), 7.1, 1.35, 5.7)
notes(s, "On the left, candidate counts per chromosome scale with chromosome size, which is what we "
      "expect for a genome-wide somatic set — no single chromosome dominates abnormally. chrY has "
      "almost nothing, consistent with HCC1395 being a female line. On the right, VAF versus depth: "
      "most candidates cluster at moderate depth across the whole VAF range, with no strange coupling "
      "between depth and VAF. In short, the spatial and joint distributions look healthy.")

# ---------- Slide 9: interpretation ----------
s = add_slide(); title(s, "Interpretation — what we can and cannot say")
bullets(s, [
    (0, "CAN say:", True),
    (1, "ClairS-TO gives a large, good-quality tumor-only candidate set (48,819 SNVs).", False),
    (1, "Coverage and distributions are well-behaved -> technically usable as markers.", False),
    (0, "CANNOT say:", True),
    (1, "that these are confirmed somatic mutations.", False),
    (1, "Tumor-only means we cannot separate true somatic from germline/artifact per variant.", False),
    (1, "The VAF ~0.5 / ~1.0 groups are germline-LIKE, not proven germline.", False),
], width=12.2, size=15)
notes(s, "This is the honest read. What we can conclude: the caller produces a large, clean-looking "
      "candidate set at good coverage, and it is technically ready to use as a marker list. What we "
      "cannot conclude: that each variant is truly somatic. Without a matched normal we cannot, "
      "variant by variant, separate somatic from germline or artifact. So I am careful not to call "
      "the 0.5 and 1.0 VAF groups germline — they merely look germline-like and would need a normal "
      "or population comparison to resolve.")

# ---------- Slide 10: limitations ----------
s = add_slide(); title(s, "Limitations")
bullets(s, [
    (0, "Tumor-only; no matched-normal validation (by design).", True),
    (0, "Possible residual germline contamination (PoN is not exhaustive).", True),
    (0, "Possible ONT sequencing / alignment artifacts.", True),
    (0, "Min-AF 0.05 removes sub-clonal low-VAF variants (clonal bias).", True),
    (0, "SNV-only — indels not characterized this week.", True),
    (0, "Candidate set NOT yet tested in the dilution series.", True),
    (0, "Context: an orthogonal check vs SEQC2 truth suggested elevated false positives —", True),
    (1, "treat absolute somatic counts cautiously.", False),
], width=12.4, size=15)
notes(s, "The limitations are mostly the direct cost of being tumor-only. There is no normal to "
      "validate against, so residual germline is expected because population databases are not "
      "complete. ONT can add its own systematic errors. The 0.05 AF floor means we miss sub-clonal "
      "variants, biasing us toward clonal markers — acceptable for MRD, but worth stating. Indels "
      "were not covered. And crucially, this marker set has not yet been looked at in the diluted "
      "samples. I also note, for honesty, that a separate comparison against SEQC2 truth hinted at a "
      "high false-positive rate, so we should not over-trust the absolute counts.")

# ---------- Slide 11: next steps / summary ----------
s = add_slide(); title(s, "Conclusion & next step")
bullets(s, [
    (0, "Bottom line: 48,819 tumor-only candidate PASS SNVs,", True),
    (1, "good coverage, well-behaved distributions,", False),
    (1, "technically ready to interrogate in the dilution BAMs.", False),
    (0, "Next research question:", True),
    (1, "Do mutant-supporting reads at these candidate loci show a consistent", False),
    (1, "dilution-dependent trend across the 1% / 0.1% / 0.01% mixed BAMs?", False),
    (0, "That test tells us whether the candidate set carries real, dose-dependent signal.", False),
], width=12.4, size=15)
notes(s, "To wrap up: we have a usable tumor-only candidate set of about 49 thousand SNVs, at good "
      "coverage, with normal-looking distributions, and it is ready to be used as predefined markers. "
      "The natural next experiment is a sanity check on signal: at these candidate loci, do the "
      "mutant-supporting reads decrease in a consistent, dose-dependent way as we go from 1% to 0.1% "
      "to 0.01% tumor fraction? If yes, the markers carry real signal and we can move toward the "
      "actual MRD detection method. If not, we revisit the candidate set first.")

prs.save(OUT)
print("saved", OUT, "(", len(prs.slides.__iter__.__self__._sldIdLst), "slides )")
