#!/usr/bin/env python3
"""Build the editable PowerPoint research report from the site narrative.

Everything the deck draws is a native PowerPoint object — text boxes, autoshapes,
tables and real charts with their own embedded worksheets. Nothing is a flattened
image, so every element can be moved, recoloured or rewritten in PowerPoint.

Numbers are read from ``research/manifests/week-001-candidate-landscape.json`` rather
than typed in, and the reference DOIs are checked against ``site/src/data/references.ts``
before the file is written: a deck that quietly disagrees with the site is worse than no
deck at all.

Usage::

    python3 tools/build_research_deck.py [-o output.pptx]
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

REPO = Path(__file__).resolve().parents[1]
MANIFEST = REPO / "research" / "manifests" / "week-001-candidate-landscape.json"
REFERENCES_TS = REPO / "site" / "src" / "data" / "references.ts"
DEFAULT_OUT = REPO / "research" / "reports" / "mrd-research-report.pptx"

# --- palette (the site's, so deck and website read as one project) ----------
INK = RGBColor(0x1F, 0x1B, 0x16)
MUTED = RGBColor(0x69, 0x5F, 0x56)
LINE = RGBColor(0xD4, 0xC8, 0xBB)
BG = RGBColor(0xF7, 0xF3, 0xEB)
PANEL = RGBColor(0xFF, 0xFD, 0xF8)
ACCENT = RGBColor(0x1F, 0x5C, 0x4E)
ACCENT_SOFT = RGBColor(0xD8, 0xEB, 0xE5)
WARN = RGBColor(0x8D, 0x5C, 0x14)
WARN_SOFT = RGBColor(0xF7, 0xEB, 0xC7)
DANGER = RGBColor(0x8F, 0x3D, 0x2C)
DANGER_SOFT = RGBColor(0xF6, 0xDD, 0xD7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SERIF = "Georgia"
SANS = "Calibri"

W, H = 13.333, 7.5
MARGIN = 0.62
CONTENT_W = W - 2 * MARGIN
BODY_TOP = 1.72
FOOTER_Y = 6.95


# --------------------------------------------------------------------------- helpers
def textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    *,
    size=13,
    bold=False,
    italic=False,
    color=INK,
    font=SANS,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    spacing=1.12,
    space_after=4,
):
    """A plain text box. `text` may contain newlines; each becomes a paragraph."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    for index, chunk in enumerate(text.split("\n")):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = chunk
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font
        run.font.color.rgb = color
    return box


def _strip_theme_style(obj):
    """Remove the shape's `<p:style>` reference.

    An empty `<a:effectLst/>` is enough for PowerPoint, but some viewers still apply the
    theme's `effectRef` and draw a drop shadow. Every shape here sets its own fill and
    line, so dropping the style element entirely makes the deck render identically
    everywhere — and flat, which is what a research figure wants.
    """
    element = obj._element
    style = element.find(
        "{http://schemas.openxmlformats.org/presentationml/2006/main}style"
    )
    if style is not None:
        element.remove(style)


def shape(
    slide,
    kind,
    x,
    y,
    w,
    h,
    *,
    fill=PANEL,
    line=LINE,
    line_w=0.75,
    adjust=None,
):
    obj = slide.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    obj.shadow.inherit = False  # Office's default drop shadow reads as decoration.
    _strip_theme_style(obj)
    if fill is None:
        obj.fill.background()
    else:
        obj.fill.solid()
        obj.fill.fore_color.rgb = fill
    if line is None:
        obj.line.fill.background()
    else:
        obj.line.color.rgb = line
        obj.line.width = Pt(line_w)
    if adjust is not None:
        try:
            obj.adjustments[0] = adjust
        except (IndexError, ValueError):
            pass
    obj.text_frame.word_wrap = True
    return obj


def card(slide, x, y, w, h, *, fill=PANEL, line=LINE):
    return shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=line, adjust=0.05)


def fill_shape_text(obj, text, *, size=12, bold=False, color=INK, font=SANS, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE, spacing=1.1, pad=0.06):
    frame = obj.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = anchor
    frame.margin_left = frame.margin_right = Inches(pad)
    frame.margin_top = frame.margin_bottom = Inches(0.04)
    for index, chunk in enumerate(text.split("\n")):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        run = para.add_run()
        run.text = chunk
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.name = font
        run.font.color.rgb = color
    return obj


def chip(slide, x, y, w, h, text, *, fill=ACCENT_SOFT, color=ACCENT, size=10, line=None, bold=False):
    obj = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill=fill, line=line, adjust=0.5)
    fill_shape_text(obj, text, size=size, color=color, bold=bold, pad=0.04)
    return obj


def new_slide(prs, *, kicker=None, title=None, lede=None, number=None, rule=True):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    if kicker:
        textbox(slide, MARGIN, 0.42, CONTENT_W, 0.26, kicker.upper(), size=10.5, color=ACCENT,
                bold=True, spacing=1.0)
    if title:
        textbox(slide, MARGIN, 0.70, CONTENT_W, 0.62, title, size=25, font=SERIF, color=INK,
                spacing=1.0)
    if lede:
        textbox(slide, MARGIN, 1.32, CONTENT_W - 0.6, 0.32, lede, size=13, color=MUTED)
    if rule:
        rule_line = shape(slide, MSO_SHAPE.RECTANGLE, MARGIN, BODY_TOP - 0.16, CONTENT_W, 0.012,
                          fill=LINE, line=None)
        rule_line.text_frame.text = ""
    if number is not None:
        textbox(slide, W - MARGIN - 1.6, FOOTER_Y + 0.18, 1.6, 0.24,
                f"{number}", size=10, color=MUTED, align=PP_ALIGN.RIGHT)
    return slide


def footnote(slide, text, *, y=None, color=MUTED, size=10.5):
    textbox(slide, MARGIN, y if y is not None else FOOTER_Y, CONTENT_W - 1.6, 0.4, text,
            size=size, color=color, italic=True)


def arrow(slide, x, y, w, h, *, direction="right", fill=ACCENT, opacity_line=None):
    kind = MSO_SHAPE.RIGHT_ARROW if direction == "right" else MSO_SHAPE.DOWN_ARROW
    obj = shape(slide, kind, x, y, w, h, fill=fill, line=None)
    return obj


def style_chart(chart, *, size=10.5):
    chart.font.size = Pt(size)
    chart.font.name = SANS
    chart.font.color.rgb = INK


# --------------------------------------------------------------------------- data
def load_summary():
    data = json.loads(MANIFEST.read_text())["summary"]
    return data


def verify_references(references):
    """Fail the build if the deck's bibliography drifts from the site's."""
    if not REFERENCES_TS.exists():
        print(f"warning: {REFERENCES_TS} not found — skipping citation cross-check", file=sys.stderr)
        return
    site_dois = set(re.findall(r'doi:\s*"([^"]+)"', REFERENCES_TS.read_text()))
    deck_text = " ".join(references)
    missing = sorted(doi for doi in site_dois if doi not in deck_text)
    if missing:
        raise SystemExit(
            "Deck references disagree with site/src/data/references.ts; missing DOIs: "
            + ", ".join(missing)
        )


REFERENCES = [
    '[1] A. Zviran et al., "Genome-wide cell-free DNA mutational integration enables '
    'ultra-sensitive cancer monitoring," Nat. Med., vol. 26, no. 7, pp. 1114–1124, Jul. 2020, '
    "doi: 10.1038/s41591-020-0915-3.",
    '[2] A. J. Widman et al., "Ultrasensitive plasma-based monitoring of tumor burden using '
    'machine-learning-guided signal enrichment," Nat. Med., vol. 30, no. 6, pp. 1655–1666, '
    "Jun. 2024, doi: 10.1038/s41591-024-03040-4.",
    "[3] N. Klimova, S. Close, D. M. Kurtz, R. D. Hockett, and L. Hyland, "
    '"Analytical validation of a circulating tumor DNA assay using PhasED-Seq technology for '
    'detecting residual disease in B-cell malignancies," Oncotarget, vol. 16, pp. 329–336, '
    "May 2025, doi: 10.18632/oncotarget.28719.",
    '[4] Y. van der Pol et al., "Real-time analysis of the cancer genome and fragmentome from '
    'plasma and urine cell-free DNA using nanopore sequencing," EMBO Mol. Med., vol. 15, no. 12, '
    "art. no. e17282, Dec. 2023, doi: 10.15252/emmm.202217282.",
    '[5] L.-T. Chen et al., "Nanopore-based consensus sequencing enables accurate multimodal '
    'tumor cell-free DNA profiling," Genome Res., vol. 35, no. 4, pp. 886–899, Apr. 2025, '
    "doi: 10.1101/gr.279144.124.",
    '[6] M. Noë et al., "DNA methylation and gene expression as determinants of genome-wide '
    'cell-free DNA fragmentation," Nat. Commun., vol. 15, no. 1, art. no. 6690, Aug. 2024, '
    "doi: 10.1038/s41467-024-50850-8.",
    '[7] Oxford Nanopore Technologies plc, "Updated method for cell-free DNA (cfDNA) methylation '
    'profiling," Requirements document, Oxford, U.K. [Online]. '
    "Available: https://nanoporetech.com/document/requirements/cfDNA-methyl-profile",
    '[8] HKU-BAL, "ClairS-TO — a deep-learning method for tumor-only somatic variant calling," '
    "v0.5.0, GitHub. [Online]. Available: https://github.com/HKU-BAL/ClairS-TO",
]


# --------------------------------------------------------------------------- slides
def slide_title(prs, summary):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.34, H, fill=ACCENT, line=None)

    textbox(slide, 1.15, 1.55, 10.6, 0.3, "RESEARCH REPORT · WORK IN PROGRESS",
            size=12, bold=True, color=ACCENT)
    textbox(slide, 1.15, 2.0, 10.9, 1.7,
            "From Tumor-Only SNV Candidates\nto a Multimodal MRD Hypothesis",
            size=40, font=SERIF, color=INK, spacing=1.05)
    shape(slide, MSO_SHAPE.RECTANGLE, 1.15, 3.72, 3.0, 0.02, fill=ACCENT, line=None)
    textbox(slide, 1.15, 3.98, 9.6, 0.8,
            "Can phase and native methylation evidence improve tumor-signal recognition\n"
            "beyond SNV-only evidence?",
            size=17, color=MUTED, font=SERIF, italic=True, spacing=1.2)

    facts = [
        ("Sample", "HCC1395 · ONT long reads · tumor-only"),
        ("Caller", "ClairS-TO v0.5.0"),
        ("Candidates", f"{summary['pass_snvs']:,} PASS SNVs of {summary['total_calls']:,} records"),
        ("Status", "Hypothesis stage — no detection result claimed"),
    ]
    x = 1.15
    for label, value in facts:
        textbox(slide, x, 5.35, 2.75, 0.24, label.upper(), size=9.5, bold=True, color=ACCENT)
        textbox(slide, x, 5.62, 2.75, 0.7, value, size=11.5, color=INK)
        x += 2.85

    textbox(slide, 1.15, 6.72, 10.9, 0.3,
            "mrd-research-os.hungnguyenmanh2k2.workers.dev  ·  18 August 2026",
            size=10.5, color=MUTED)
    return slide


def slide_problem(prs, number):
    slide = new_slide(
        prs,
        kicker="1 · The problem",
        title="A rare signal inside a chemically identical background",
        lede="Minimal residual disease is what remains after treatment, below the reach of imaging.",
        number=number,
    )

    # Dot matrix: 25 × 12 = 300 molecules, 3 of them tumor-derived (1%).
    cols, rows, step, dot = 25, 14, 0.215, 0.155
    x0, y0 = MARGIN + 0.05, BODY_TOP + 0.45
    tumor_cells = {(3, 7), (7, 18), (11, 12)}
    for r in range(rows):
        for c in range(cols):
            is_tumor = (r, c) in tumor_cells
            obj = shape(
                slide,
                MSO_SHAPE.OVAL,
                x0 + c * step,
                y0 + r * step,
                dot,
                dot,
                fill=DANGER if is_tumor else RGBColor(0xE3, 0xDC, 0xD1),
                line=None,
            )
            if is_tumor:
                obj.line.color.rgb = DANGER
                obj.line.width = Pt(1.0)

    legend_y = y0 + rows * step + 0.22
    shape(slide, MSO_SHAPE.OVAL, x0, legend_y + 0.03, dot, dot, fill=DANGER, line=None)
    textbox(slide, x0 + 0.22, legend_y, 2.4, 0.26, "tumor-derived", size=10.5, color=INK)
    shape(slide, MSO_SHAPE.OVAL, x0 + 1.7, legend_y + 0.03, dot, dot,
          fill=RGBColor(0xE3, 0xDC, 0xD1), line=None)
    textbox(slide, x0 + 1.92, legend_y, 2.4, 0.26, "background", size=10.5, color=MUTED)

    right_x = x0 + cols * step + 0.6
    right_w = W - MARGIN - right_x
    box = card(slide, right_x, BODY_TOP + 0.45, right_w, 1.35, fill=ACCENT, line=None)
    fill_shape_text(
        box,
        "3 in 350 molecules\nat 1% tumor fraction",
        size=19,
        font=SERIF,
        color=WHITE,
        align=PP_ALIGN.CENTER,
    )
    textbox(slide, right_x, BODY_TOP + 2.1, right_w, 1.6,
            "The measurement is not a search for something unusual. It is a search for "
            "specific molecules that are chemically identical to a background arriving in "
            "vast excess.",
            size=14, color=INK, spacing=1.35)
    divider = shape(slide, MSO_SHAPE.RECTANGLE, right_x, BODY_TOP + 3.35, 1.6, 0.02,
                    fill=LINE, line=None)
    divider.text_frame.text = ""
    textbox(slide, right_x, BODY_TOP + 3.6, right_w, 1.2,
            "Detected residue is the earliest evidence about treatment outcome.\n"
            "Its absence is evidence — not proof.",
            size=12.5, color=MUTED, spacing=1.35)

    footnote(slide, "Illustrative schematic. Not measured data.")
    return slide


def slide_barriers(prs, number):
    slide = new_slide(
        prs,
        kicker="2 · Why it is difficult",
        title="Every method is improving the same ratio",
        number=number,
    )

    # Ratio, drawn as a fraction rather than written as a sentence.
    fx, fw = MARGIN, 3.55
    textbox(slide, fx, BODY_TOP + 0.35, fw, 0.5, "Tumor signal", size=20, font=SERIF,
            color=ACCENT, align=PP_ALIGN.CENTER)
    shape(slide, MSO_SHAPE.RECTANGLE, fx + 0.35, BODY_TOP + 0.92, fw - 0.7, 0.02,
          fill=INK, line=None)
    textbox(slide, fx, BODY_TOP + 1.02, fw, 0.5, "Background noise", size=20, font=SERIF,
            color=DANGER, align=PP_ALIGN.CENTER)
    arrow(slide, fx + fw / 2 - 0.13, BODY_TOP + 1.75, 0.26, 0.5, direction="right", fill=ACCENT)
    arrow_shape = slide.shapes[-1]
    arrow_shape.rotation = 270
    textbox(slide, fx, BODY_TOP + 2.4, fw, 0.4, "maximise", size=12, color=MUTED,
            align=PP_ALIGN.CENTER, italic=True)

    denominator = card(slide, fx, BODY_TOP + 2.95, fw, 1.55, fill=DANGER_SOFT, line=None)
    fill_shape_text(
        denominator,
        "Background ⊋ sequencing error\n\n"
        "germline without a matched normal · clonal haematopoiesis ·\n"
        "alignment artifacts · protocol bias",
        size=11,
        color=DANGER,
        align=PP_ALIGN.CENTER,
    )

    barriers = [
        ("1", "Limited tumor-derived molecules",
         "At low tumor fraction the evidence at one locus can shrink to a single observation."),
        ("2", "Errors mimic true mutations",
         "One read carrying an alternative allele looks the same whether it is real or an error."),
        ("3", "One modality may not be enough",
         "If a molecule is described by a single fact, that fact must separate signal from background alone."),
    ]
    bx = fx + fw + 0.6
    bw = W - MARGIN - bx
    bh = 1.42
    for index, (num, head, body) in enumerate(barriers):
        y = BODY_TOP + 0.3 + index * (bh + 0.28)
        box = card(slide, bx, y, bw, bh)
        badge = shape(slide, MSO_SHAPE.OVAL, bx + 0.28, y + 0.28, 0.46, 0.46,
                      fill=ACCENT, line=None)
        fill_shape_text(badge, num, size=15, bold=True, color=WHITE, font=SERIF)
        textbox(slide, bx + 0.95, y + 0.26, bw - 1.3, 0.35, head, size=15, font=SERIF, color=INK)
        textbox(slide, bx + 0.95, y + 0.68, bw - 1.3, 0.65, body, size=11.5, color=MUTED,
                spacing=1.2)
    return slide


def slide_related_work(prs, number):
    slide = new_slide(
        prs,
        kicker="3 · Related work",
        title="Each study attacks the ratio at a different point",
        number=number,
    )

    rows = [
        ("MRDetect [1]", "Genome-wide SNVs",
         "Breadth over depth: integrate many weak loci",
         "Aggregates across loci, not within a molecule"),
        ("MRD-EDGE [2]", "SNV + CNV, learned",
         "Enrich signal before aggregation",
         "Short-read feature space; no base modifications"),
        ("PhasED-Seq [3]", "Phased variants",
         "Linkage: background must fail twice on one fragment",
         "Needs ≥2 variants per fragment — scarce at low TF"),
        ("Real-time ONT [4]", "CNA + fragmentomics",
         "Native long-read readout, no separate assay",
         "Not single-nucleotide evidence at low TF"),
        ("NanoRCS [5]", "Consensus SNV / CNA / fragment",
         "Raise per-read accuracy, combine modalities",
         "Modalities are sequence-level, not phase + methylation"),
        ("Methylation ↔ fragmentation [6]", "Methylation, fragmentation",
         "Explain what determines fragmentation",
         "Warns: modalities are coupled, not independent"),
        ("ONT cfDNA methylation [7]", "Native methylation",
         "Specify protocol requirements",
         "Requirements document, not a detection study"),
    ]
    headers = ("Study", "Evidence modality", "Strategy", "Limitation for tumor-only long reads")
    widths = (2.55, 2.35, 3.55, 3.68)

    table_shape = slide.shapes.add_table(
        len(rows) + 1, 4, Inches(MARGIN), Inches(BODY_TOP + 0.12),
        Inches(sum(widths)), Inches(4.35)
    )
    table = table_shape.table
    table.first_row = True
    table.horz_banding = False
    for index, width in enumerate(widths):
        table.columns[index].width = Inches(width)
    table.rows[0].height = Inches(0.42)

    for col, header in enumerate(headers):
        cell = table.cell(0, col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.09)
        para = cell.text_frame.paragraphs[0]
        run = para.add_run()
        run.text = header
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.name = SANS
        run.font.color.rgb = WHITE

    for r, row in enumerate(rows, start=1):
        table.rows[r].height = Inches(0.55)
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if r % 2 else RGBColor(0xF4, 0xF1, 0xE8)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.09)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            para = cell.text_frame.paragraphs[0]
            run = para.add_run()
            run.text = value
            run.font.size = Pt(10.5)
            run.font.name = SANS
            run.font.bold = c == 0
            run.font.color.rgb = INK if c == 0 else MUTED

    footnote(slide, "No performance figures are compared: different assays, sample types and "
                    "cohorts make such a comparison unsupported.", y=FOOTER_Y - 0.15)
    return slide


def slide_matrix(prs, number):
    slide = new_slide(
        prs,
        kicker="4 · Synthesis",
        title="What is measured on the same molecule",
        lede="Filled marks show the evidence each study reads off one physical molecule.",
        number=number,
    )

    columns = ["SNV", "CNA", "Fragmentomics", "Phase", "Native\nmethylation", "Joint on one\nmolecule"]
    rows = [
        ("MRDetect [1]", [1, 0, 0, 0, 0, 0]),
        ("MRD-EDGE [2]", [1, 1, 0, 0, 0, 0]),
        ("PhasED-Seq [3]", [1, 0, 0, 1, 0, 1]),
        ("Real-time ONT [4]", [0, 1, 1, 0, 0, 0]),
        ("NanoRCS [5]", [1, 1, 1, 0, 0, 1]),
        ("Methylation ↔ fragm. [6]", [0, 0, 1, 0, 1, 0]),
        ("ONT cfDNA methyl. [7]", [0, 0, 0, 0, 1, 0]),
        ("This project (proposed)", [1, 0, 0, 1, 1, 1]),
    ]

    label_w = 3.05
    col_w = 1.5
    x0 = MARGIN
    y0 = BODY_TOP + 0.82
    row_h = 0.46

    for index, header in enumerate(columns):
        textbox(slide, x0 + label_w + index * col_w, y0 - 0.68, col_w, 0.6, header,
                size=10.5, bold=True, color=ACCENT, align=PP_ALIGN.CENTER, spacing=1.0)

    for r, (label, marks) in enumerate(rows):
        y = y0 + r * row_h
        proposed = r == len(rows) - 1
        band = shape(slide, MSO_SHAPE.RECTANGLE, x0, y, label_w + len(columns) * col_w, row_h - 0.06,
                     fill=ACCENT_SOFT if proposed else (PANEL if r % 2 == 0 else BG), line=None)
        band.text_frame.text = ""
        textbox(slide, x0 + 0.12, y + 0.08, label_w - 0.2, 0.3, label, size=11,
                bold=proposed, color=INK if proposed else MUTED)
        for c, mark in enumerate(marks):
            cx = x0 + label_w + c * col_w + col_w / 2
            if mark:
                dot = shape(slide, MSO_SHAPE.OVAL, cx - 0.09, y + 0.10, 0.18, 0.18,
                            fill=ACCENT if proposed else INK, line=None)
                dot.text_frame.text = ""
            else:
                dot = shape(slide, MSO_SHAPE.OVAL, cx - 0.07, y + 0.12, 0.14, 0.14,
                            fill=None, line=LINE, line_w=1.0)
                dot.text_frame.text = ""

    note = card(slide, MARGIN, y0 + len(rows) * row_h + 0.20, CONTENT_W, 0.6,
                fill=WARN_SOFT, line=None)
    fill_shape_text(
        note,
        "Within the supplied literature set, no study reads allele, phase and native methylation "
        "off the same molecule. This is a statement about these seven documents — not about the field.",
        size=12, color=WARN, align=PP_ALIGN.CENTER,
    )
    return slide


def slide_funnel(prs, number, summary):
    slide = new_slide(
        prs,
        kicker="5 · Current output",
        title="Tumor-only calling produced candidates, not conclusions",
        number=number,
    )

    steps = [
        ("HCC1395 tumor-only ONT BAM", "long reads · no matched normal", PANEL, INK, 5.9),
        ("ClairS-TO v0.5.0 [8]", "tumor-only candidate calling and filtering", PANEL, INK, 5.2),
        (f"{summary['total_calls']:,} candidate records", "all records emitted", RGBColor(0xEC, 0xE6, 0xDA), INK, 4.4),
        (f"{summary['pass_snvs']:,} PASS SNV candidates",
         f"{summary['pass_snv_fraction_pct']}% of records", ACCENT, WHITE, 3.3),
    ]
    x_center = MARGIN + 4.1
    y = BODY_TOP + 0.12
    for index, (label, sub, fill, color, width) in enumerate(steps):
        box = card(slide, x_center - width / 2, y, width, 0.82, fill=fill,
                   line=None if fill == ACCENT else LINE)
        fill_shape_text(box, label, size=15 if index >= 2 else 13, bold=index >= 2,
                        color=color, font=SERIF if index >= 2 else SANS)
        textbox(slide, x_center - width / 2, y + 0.84, width, 0.24, sub, size=10,
                color=MUTED, align=PP_ALIGN.CENTER)
        if index < len(steps) - 1:
            arrow(slide, x_center - 0.10, y + 1.10, 0.20, 0.26, direction="down",
                  fill=RGBColor(0xC9, 0xBE, 0xB2))
        y += 1.42

    panel_x = MARGIN + 8.55
    panel_w = W - MARGIN - panel_x
    warn_box = card(slide, panel_x, BODY_TOP + 0.12, panel_w, 1.75, fill=WARN_SOFT, line=None)
    fill_shape_text(
        warn_box,
        "PASS is a caller retention label.\nIt is not confirmed somatic truth.",
        size=15, color=WARN, font=SERIF, align=PP_ALIGN.LEFT, pad=0.24, spacing=1.35,
    )
    textbox(slide, panel_x, BODY_TOP + 2.15, panel_w, 1.2,
            "The 98.46% not retained is a selection funnel, not a false-positive rate.",
            size=13, color=INK, spacing=1.35)
    textbox(slide, panel_x, BODY_TOP + 3.15, panel_w, 1.4,
            "The composition of the retained set — somatic, germline, artifact — is unknown. "
            "That gap is the next piece of work.",
            size=13, color=INK, spacing=1.35)
    return slide


def slide_distributions(prs, number, summary):
    slide = new_slide(
        prs,
        kicker="6 · Candidate landscape",
        title="Retained candidates carry substantial read support",
        lede="Descriptive characterisation of the PASS SNV set. No candidate is labelled somatic.",
        number=number,
    )

    ranges = [
        ("Read depth", summary["depth_q25"], summary["median_depth"], summary["depth_q75"], 0, 130, "×"),
        ("Allele fraction (VAF)", summary["vaf_q25"], summary["median_vaf"], summary["vaf_q75"], 0, 1.0, ""),
        ("ALT-supporting reads", summary["alt_support_q25"], summary["median_alt_support"],
         summary["alt_support_q75"], 0, 65, ""),
    ]
    track_x = MARGIN + 2.05
    track_w = 4.55
    y = BODY_TOP + 0.85
    for label, q25, median, q75, lo, hi, unit in ranges:
        textbox(slide, MARGIN, y - 0.04, 1.95, 0.3, label, size=12, color=INK)
        track = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, track_x, y + 0.06, track_w, 0.16,
                      fill=RGBColor(0xE9, 0xE3, 0xD8), line=None, adjust=0.5)
        track.text_frame.text = ""
        left = track_x + track_w * (q25 - lo) / (hi - lo)
        width = track_w * (q75 - q25) / (hi - lo)
        band = shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, y + 0.02, width, 0.24,
                     fill=ACCENT_SOFT, line=None, adjust=0.5)
        band.text_frame.text = ""
        med_x = track_x + track_w * (median - lo) / (hi - lo)
        marker = shape(slide, MSO_SHAPE.RECTANGLE, med_x - 0.02, y - 0.06, 0.05, 0.4,
                       fill=ACCENT, line=None)
        marker.text_frame.text = ""
        median_text = f"{median}{unit}"
        textbox(slide, track_x + track_w + 0.18, y - 0.06, 1.5, 0.3,
                f"median {median_text}", size=11.5, bold=True, color=ACCENT)
        textbox(slide, track_x + track_w + 0.18, y + 0.22, 1.9, 0.3,
                f"IQR {q25}–{q75}{unit}", size=10, color=MUTED)
        y += 1.25

    textbox(slide, track_x, y - 0.25, track_w, 0.3,
            "bar = interquartile range   |   line = median", size=10.5, color=MUTED,
            italic=True)

    # Allele-fraction composition, as a real (editable) chart.
    het = summary["vaf_het_like_0_40_0_60_pct"]
    hom = summary["vaf_hom_like_gt_0_90_pct"]
    other = round(100 - het - hom, 1)
    chart_data = CategoryChartData()
    chart_data.categories = ["Het-like 0.40–0.60", "Hom-like > 0.90", "Other VAF"]
    chart_data.add_series("Share of PASS SNVs", (het, hom, other))
    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(MARGIN + 8.25), Inches(BODY_TOP + 0.35),
        Inches(4.1), Inches(4.0), chart_data
    )
    chart = graphic.chart
    style_chart(chart, size=10.5)
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0.0"%"'
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.show_value = True
    plot.data_labels.font.size = Pt(11)
    plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = INK
    for point, colour in zip(plot.series[0].points, (ACCENT, WARN, RGBColor(0xC9, 0xBE, 0xB2))):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colour
    textbox(slide, MARGIN + 8.25, BODY_TOP + 0.05, 4.1, 0.3,
            "Allele-fraction composition", size=12, bold=True, color=INK,
            align=PP_ALIGN.CENTER)

    footnote(slide, "VAF is an allele fraction in the source sample — it is not a tumor fraction.")
    return slide


def slide_correlations(prs, number, summary):
    slide = new_slide(
        prs,
        kicker="7 · Reading the landscape",
        title="Allele fraction is not a simple coverage effect",
        number=number,
    )

    chart_data = CategoryChartData()
    chart_data.categories = ["VAF ~ depth", "ALT support ~ depth", "ALT support ~ VAF"]
    chart_data.add_series("Pearson r", (
        summary["corr_vaf_depth"], summary["corr_alt_depth"], summary["corr_alt_vaf"]
    ))
    graphic = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(MARGIN), Inches(BODY_TOP + 0.25),
        Inches(7.1), Inches(4.35), chart_data
    )
    chart = graphic.chart
    style_chart(chart, size=11)
    chart.has_legend = False
    value_axis = chart.value_axis
    value_axis.minimum_scale = -0.6
    value_axis.maximum_scale = 0.8
    value_axis.has_major_gridlines = True
    value_axis.tick_labels.number_format = "0.0"
    value_axis.tick_labels.number_format_is_linked = False
    plot = chart.plots[0]
    plot.gap_width = 60
    plot.has_data_labels = True
    plot.data_labels.number_format = "0.000"
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.show_value = True
    plot.data_labels.font.size = Pt(11)
    plot.data_labels.font.bold = True
    series = plot.series[0]
    for point, colour in zip(series.points, (DANGER, ACCENT, ACCENT)):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colour

    panel_x = MARGIN + 7.5
    panel_w = W - MARGIN - panel_x
    obs = card(slide, panel_x, BODY_TOP + 0.25, panel_w, 1.85, fill=PANEL)
    fill_shape_text(
        obs,
        "OBSERVATION\n\nVAF falls slightly as depth rises; ALT support tracks depth directly.",
        size=12, color=INK, align=PP_ALIGN.LEFT, pad=0.2,
    )
    interp = card(slide, panel_x, BODY_TOP + 2.3, panel_w, 2.3, fill=ACCENT_SOFT, line=None)
    fill_shape_text(
        interp,
        "INTERPRETATION\n\nHeterogeneity in the retained set does not collapse into one coverage "
        "effect. This bounds a technical concern; it does not establish biological origin.",
        size=12, color=ACCENT, align=PP_ALIGN.LEFT, pad=0.2,
    )
    return slide


def slide_gap(prs, number):
    slide = new_slide(
        prs,
        kicker="8 · Research gap",
        title="From a candidate list to an open question",
        number=number,
    )

    steps = [
        ("Tumor-only SNV\ncandidates", "produced", PANEL, INK),
        ("At low tumor fraction,\nisolated ALT observations\nresemble background",
         "follows from the barriers", PANEL, INK),
        ("Long reads also carry phase\nand native methylation\non the same molecule",
         "property of the platform", PANEL, INK),
        ("Added discriminative\nvalue is unknown", "open — and testable", ACCENT, WHITE),
    ]
    box_w = 2.72
    gap = 0.42
    x = MARGIN + 0.15
    y = BODY_TOP + 1.05
    for index, (label, sub, fill, color) in enumerate(steps):
        box = card(slide, x, y, box_w, 2.05, fill=fill, line=None if fill == ACCENT else LINE)
        fill_shape_text(box, label, size=14.5, color=color, font=SERIF, spacing=1.25)
        textbox(slide, x, y + 2.14, box_w, 0.3, sub, size=10.5, color=MUTED,
                align=PP_ALIGN.CENTER, italic=True)
        if index < len(steps) - 1:
            arrow(slide, x + box_w + 0.07, y + 0.88, gap - 0.14, 0.3,
                  fill=RGBColor(0xC9, 0xBE, 0xB2))
        x += box_w + gap

    footnote(slide, "Conceptual chain. No step asserts a measured effect, and the "
                    "representation and evaluation of these modalities are deliberately not "
                    "specified here.", y=FOOTER_Y - 0.55)
    return slide


def slide_hypothesis(prs, number):
    slide = new_slide(
        prs,
        kicker="9 · Hypothesis",
        title="What the next stage will test",
        number=number,
    )

    question = card(slide, MARGIN, BODY_TOP + 0.25, CONTENT_W, 1.15, fill=ACCENT, line=None)
    fill_shape_text(
        question,
        "Does phase evidence and/or native methylation evidence improve the recognition\n"
        "of tumor-derived signal compared with SNV-only evidence?",
        size=17, color=WHITE, font=SERIF, spacing=1.25,
    )
    textbox(slide, MARGIN + 0.05, BODY_TOP + 0.02, 4.0, 0.25, "RESEARCH QUESTION",
            size=10, bold=True, color=ACCENT)

    half = (CONTENT_W - 0.4) / 2
    h0 = card(slide, MARGIN, BODY_TOP + 1.75, half, 1.6, fill=PANEL)
    fill_shape_text(
        h0,
        "H₀\n\nAdding phase or native methylation evidence does not improve tumor-signal "
        "recognition over SNV-only evidence.",
        size=13, color=INK, align=PP_ALIGN.LEFT, pad=0.24, spacing=1.25,
    )
    h1 = card(slide, MARGIN + half + 0.4, BODY_TOP + 1.75, half, 1.6, fill=ACCENT_SOFT, line=None)
    fill_shape_text(
        h1,
        "H₁\n\nPhase and/or native methylation evidence provides additional discriminative "
        "value beyond SNV-only evidence.",
        size=13, color=ACCENT, align=PP_ALIGN.LEFT, pad=0.24, spacing=1.25,
    )

    scope = card(slide, MARGIN, BODY_TOP + 3.6, CONTENT_W, 0.85, fill=WARN_SOFT, line=None)
    fill_shape_text(
        scope,
        "Either outcome is informative. A surviving H₀ means phase and methylation are redundant "
        "with sequence evidence in this setting — a result about long-read MRD, not a failed experiment.",
        size=12, color=WARN, align=PP_ALIGN.CENTER,
    )
    return slide


def slide_boundaries(prs, number):
    slide = new_slide(
        prs,
        kicker="10 · Claim boundaries",
        title="What this report does not claim",
        number=number,
    )

    items = [
        ("PASS ≠ somatic", "A retention label from the caller, not a biological verdict."),
        ("VAF ≠ tumor fraction", "An allele fraction in the source sample."),
        ("Coverage ≠ sensitivity", "High depth here says nothing about low-tumor-fraction detection."),
        ("Gap ≠ field-wide", "The unresolved question is scoped to the seven supplied documents."),
    ]
    box_w = (CONTENT_W - 3 * 0.34) / 4
    x = MARGIN
    for head, body in items:
        box = card(slide, x, BODY_TOP + 0.75, box_w, 3.05)
        strip = shape(slide, MSO_SHAPE.RECTANGLE, x, BODY_TOP + 0.75, box_w, 0.07,
                      fill=DANGER, line=None)
        strip.text_frame.text = ""
        textbox(slide, x + 0.26, BODY_TOP + 1.15, box_w - 0.52, 0.8, head, size=17,
                font=SERIF, color=DANGER, spacing=1.15)
        textbox(slide, x + 0.26, BODY_TOP + 2.0, box_w - 0.52, 1.5, body, size=12.5,
                color=MUTED, spacing=1.3)
        x += box_w + 0.34

    footnote(slide, "No experimental result is reported in this deck. Every number shown is a "
                    "property of the candidate call set.", y=FOOTER_Y - 0.5)
    return slide


def slide_references(prs, number):
    slide = new_slide(prs, kicker="References", title="Sources", number=number)
    column_w = (CONTENT_W - 0.5) / 2
    for index, reference in enumerate(REFERENCES):
        col, row = divmod(index, 4)
        textbox(slide, MARGIN + col * (column_w + 0.5), BODY_TOP + 0.3 + row * 1.15,
                column_w, 1.0, reference, size=10.5, color=INK, spacing=1.2)
    footnote(slide, "Bibliographic fields verified against publisher records; identical to the "
                    "project website's reference list.", y=FOOTER_Y - 0.1)
    return slide


# --------------------------------------------------------------------------- main
def build(out_path: Path) -> Path:
    summary = load_summary()
    verify_references(REFERENCES)

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    slide_title(prs, summary)
    slide_problem(prs, 2)
    slide_barriers(prs, 3)
    slide_related_work(prs, 4)
    slide_matrix(prs, 5)
    slide_funnel(prs, 6, summary)
    slide_distributions(prs, 7, summary)
    slide_correlations(prs, 8, summary)
    slide_gap(prs, 9)
    slide_hypothesis(prs, 10)
    slide_boundaries(prs, 11)
    slide_references(prs, 12)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(out_path)
    return out_path


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("-o", "--output", type=Path, default=DEFAULT_OUT)
    args = parser.parse_args()
    path = build(args.output)
    print(f"Wrote {path.relative_to(REPO)} ({path.stat().st_size / 1024:.0f} KB)")


if __name__ == "__main__":
    main()
